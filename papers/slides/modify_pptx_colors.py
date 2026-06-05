#!/usr/bin/env python
"""
Modify PowerPoint presentation to have dark background and light text.

Example usage:
> modify_pptx_colors.py --input ref.pptx --output ref_dark.pptx
> modify_pptx_colors.py --input ref.pptx --output ref_bg.pptx --background image.png

Import as:

import modify_pptx_colors as moppcol
"""

import argparse
import logging
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

import helpers.hdbg as hdbg
import helpers.hparser as hparser

_LOG = logging.getLogger(__name__)

# #############################################################################
# Constants
# #############################################################################


# Dark background (dark gray/charcoal).
_DEFAULT_BG_COLOR = "#2D2D2D"
# Light text.
_DEFAULT_TEXT_COLOR = "#FAFBFD"


# #############################################################################
# Helper functions
# #############################################################################


def _hex_to_rgb(hex_color: str) -> tuple:
    """
    Convert hex color to RGB tuple.

    :param hex_color: hex color string (e.g., "#2D2D2D")
    :return: RGB tuple (r, g, b)
    """
    hex_color = hex_color.lstrip("#")
    hdbg.dassert_eq(
        len(hex_color),
        6,
        "Invalid hex color format. Expected 6 characters:",
        hex_color,
    )
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def _set_text_color(shape, color_rgb: tuple) -> None:
    """
    Set text color for a shape.

    :param shape: PowerPoint shape object
    :param color_rgb: RGB color tuple (r, g, b)
    """
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(*color_rgb)


def _set_slide_background(slide, color_rgb: tuple) -> None:
    """
    Set slide background color.

    :param slide: PowerPoint slide object
    :param color_rgb: RGB color tuple (r, g, b)
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color_rgb)


def _set_slide_background_image(slide, prs, image_path: str) -> None:
    """
    Set slide background image.

    :param slide: PowerPoint slide object
    :param prs: Presentation object
    :param image_path: path to background image file
    """
    # Get slide dimensions.
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    # Add picture to slide at full size starting at top-left corner.
    left = top = 0
    pic = slide.shapes.add_picture(
        image_path, left, top, width=slide_width, height=slide_height
    )
    # Move picture to back (below all other shapes).
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def _modify_presentation(
    input_file: str,
    output_file: str,
    *,
    bg_color_hex: str = _DEFAULT_BG_COLOR,
    text_color_hex: str = _DEFAULT_TEXT_COLOR,
    background_image: Optional[str] = None,
) -> None:
    """
    Modify presentation colors and background.

    :param input_file: path to input PowerPoint file
    :param output_file: path to output PowerPoint file
    :param bg_color_hex: background color in hex format
    :param text_color_hex: text color in hex format
    :param background_image: optional path to background image file
    """
    hdbg.dassert_file_exists(input_file)
    if background_image is not None:
        hdbg.dassert_file_exists(background_image)
    # Load presentation.
    _LOG.info("Loading presentation from: %s", input_file)
    prs = Presentation(input_file)
    # Convert colors.
    text_rgb = _hex_to_rgb(text_color_hex)
    _LOG.info("Using text color: %s (RGB: %s)", text_color_hex, text_rgb)
    if background_image is not None:
        _LOG.info("Using background image: %s", background_image)
    else:
        bg_rgb = _hex_to_rgb(bg_color_hex)
        _LOG.info("Using background color: %s (RGB: %s)", bg_color_hex, bg_rgb)
    # Process slides.
    _LOG.info("Processing %s slides...", len(prs.slides))
    for slide_idx, slide in enumerate(prs.slides, 1):
        _LOG.debug("Processing slide %s...", slide_idx)
        # Set background.
        if background_image is not None:
            _set_slide_background_image(slide, prs, background_image)
        else:
            _set_slide_background(slide, bg_rgb)
        # Update text colors.
        for shape in slide.shapes:
            if shape.has_text_frame:
                _set_text_color(shape, text_rgb)
            # Handle tables.
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(*text_rgb)
            # Handle groups (recursively process shapes in groups).
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.has_text_frame:
                        _set_text_color(sub_shape, text_rgb)
    # Save presentation.
    _LOG.info("Saving modified presentation to: %s", output_file)
    prs.save(output_file)
    _LOG.info("Done!")


# #############################################################################


def _parse() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=__doc__,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "--input",
        action="store",
        required=True,
        help="Input PowerPoint file path",
    )
    parser.add_argument(
        "--output",
        action="store",
        required=True,
        help="Output PowerPoint file path",
    )
    parser.add_argument(
        "--bg_color",
        action="store",
        default=_DEFAULT_BG_COLOR,
        help="Background color in hex format (default: %s)" % _DEFAULT_BG_COLOR,
    )
    parser.add_argument(
        "--text_color",
        action="store",
        default=_DEFAULT_TEXT_COLOR,
        help="Text color in hex format (default: %s)" % _DEFAULT_TEXT_COLOR,
    )
    parser.add_argument(
        "--background",
        action="store",
        default=None,
        help="Background image file path (overrides --bg_color if provided)",
    )
    hparser.add_verbosity_arg(parser)
    return parser


def _main(parser: argparse.ArgumentParser) -> None:
    args = parser.parse_args()
    hdbg.init_logger(verbosity=args.log_level, use_exec_path=True)
    _modify_presentation(
        args.input,
        args.output,
        bg_color_hex=args.bg_color,
        text_color_hex=args.text_color,
        background_image=args.background,
    )


if __name__ == "__main__":
    _main(_parse())
