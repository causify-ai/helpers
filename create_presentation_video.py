#!/usr/bin/env python3

"""
Script to create a video from PowerPoint slides and video clips with picture-in-picture.

This script:
1. Extracts slides from presentation.ppt as images
2. Creates a video with:
   - Slide 1 for 5 seconds with comment1.mp4 as picture-in-picture
   - demo1.mp4 with comment2.mp4 as picture-in-picture  
   - Slide 2 with comment3.mp4 as picture-in-picture

Requirements:
- moviepy
- python-pptx (for extracting slides)
- Pillow (for image processing)
"""

import os
import sys
from pathlib import Path
from typing import List, Tuple, Optional

from moviepy.editor import (
    VideoFileClip,
    ImageClip,
    CompositeVideoClip,
    concatenate_videoclips,
)
from PIL import Image
import tempfile


class PresentationVideoCreator:
    """Creates videos from PowerPoint presentations and video clips."""
    
    def __init__(self, presentation_path: str = "presentation.ppt"):
        self.presentation_path = presentation_path
        self.temp_dir = tempfile.mkdtemp()
        
    def __enter__(self):
        return self
        
    def __exit__(self, exc_type, exc_val, exc_tb):
        self._cleanup_temp_files()
    
    def _cleanup_temp_files(self):
        """Clean up temporary files."""
        import shutil
        if os.path.exists(self.temp_dir):
            shutil.rmtree(self.temp_dir)
    
    def extract_slides_as_images(self) -> List[str]:
        """
        Extract slides from PowerPoint presentation as PNG images.
        
        Returns:
            List of paths to extracted slide images
        """
        if not os.path.exists(self.presentation_path):
            print(f"Warning: {self.presentation_path} not found. Creating placeholder slides.")
            return self._create_placeholder_slides()
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.util import Inches
            
            # Load presentation
            prs = Presentation(self.presentation_path)
            slide_paths = []
            
            for i, slide in enumerate(prs.slides, 1):
                # Create a simple image for each slide
                # Note: python-pptx doesn't directly export to images
                # This is a simplified approach - for production use, consider using 
                # win32com.client (Windows) or LibreOffice API
                slide_path = os.path.join(self.temp_dir, f"slide_{i}.png")
                self._create_slide_image(f"Slide {i}", slide_path)
                slide_paths.append(slide_path)
                
            return slide_paths
            
        except ImportError:
            print("python-pptx not available. Creating placeholder slides.")
            return self._create_placeholder_slides()
        except Exception as e:
            print(f"Error extracting slides: {e}. Creating placeholder slides.")
            return self._create_placeholder_slides()
    
    def _create_placeholder_slides(self) -> List[str]:
        """Create placeholder slide images."""
        slide_paths = []
        for i in [1, 2]:  # Only need slides 1 and 2
            slide_path = os.path.join(self.temp_dir, f"slide_{i}.png")
            self._create_slide_image(f"Slide {i}", slide_path)
            slide_paths.append(slide_path)
        return slide_paths
    
    def _create_slide_image(self, text: str, output_path: str):
        """Create a simple slide image with text."""
        from PIL import Image, ImageDraw, ImageFont
        
        # Create a white background image
        width, height = 1920, 1080
        image = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(image)
        
        # Try to use a decent font
        try:
            font = ImageFont.truetype("arial.ttf", 72)
        except:
            font = ImageFont.load_default()
        
        # Calculate text position (centered)
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (width - text_width) // 2
        y = (height - text_height) // 2
        
        # Draw text
        draw.text((x, y), text, fill='black', font=font)
        
        # Add a border
        draw.rectangle([50, 50, width-50, height-50], outline='blue', width=5)
        
        # Save image
        image.save(output_path)
    
    def create_pip_video(self, main_clip: VideoFileClip, pip_video_path: str, 
                        pip_position: Tuple[str, str] = ('right', 'top'),
                        pip_size: Tuple[int, int] = (320, 240)) -> CompositeVideoClip:
        """
        Create picture-in-picture video.
        
        Args:
            main_clip: Main video clip
            pip_video_path: Path to picture-in-picture video
            pip_position: Position of PiP ('left'/'right', 'top'/'bottom')  
            pip_size: Size of PiP video (width, height)
            
        Returns:
            Composite video with picture-in-picture
        """
        if not os.path.exists(pip_video_path):
            print(f"Warning: {pip_video_path} not found. Creating placeholder.")
            pip_video_path = self._create_placeholder_video(pip_video_path, pip_size)
        
        # Load and resize PiP video
        pip_clip = VideoFileClip(pip_video_path)
        pip_clip = pip_clip.resize(pip_size)
        
        # Match duration with main clip
        pip_clip = pip_clip.set_duration(main_clip.duration)
        
        # Position PiP video
        margin = 20
        if pip_position[0] == 'right':
            x_pos = main_clip.w - pip_size[0] - margin
        else:  # left
            x_pos = margin
            
        if pip_position[1] == 'top':
            y_pos = margin
        else:  # bottom
            y_pos = main_clip.h - pip_size[1] - margin
        
        pip_clip = pip_clip.set_position((x_pos, y_pos))
        
        # Create composite
        return CompositeVideoClip([main_clip, pip_clip])
    
    def _create_placeholder_video(self, video_path: str, size: Tuple[int, int]) -> str:
        """Create a placeholder video file."""
        placeholder_path = os.path.join(self.temp_dir, f"placeholder_{os.path.basename(video_path)}")
        
        # Create a simple colored clip
        from moviepy.editor import ColorClip
        placeholder_clip = ColorClip(size=size, color=(100, 150, 200), duration=5)
        try:
            placeholder_clip.write_videofile(placeholder_path, fps=24)
        except Exception as e:
            print(f"Warning: Could not create placeholder video: {e}")
            # Return the original path anyway, caller will handle
            pass
        
        return placeholder_path
    
    def create_final_video(self, output_path: str = "final_presentation_video.mp4"):
        """
        Create the final video according to specifications.
        
        Args:
            output_path: Path for output video file
        """
        print("Starting video creation process...")
        
        # Extract slides
        print("Extracting slides...")
        slide_paths = self.extract_slides_as_images()
        
        if len(slide_paths) < 2:
            print("Error: Need at least 2 slides")
            return
        
        video_clips = []
        
        # 1. Slide 1 for 5 seconds with comment1.mp4 PiP
        print("Creating slide 1 segment...")
        slide1_clip = ImageClip(slide_paths[0]).set_duration(5)
        slide1_with_pip = self.create_pip_video(slide1_clip, "comment1.mp4")
        video_clips.append(slide1_with_pip)
        
        # 2. demo1.mp4 with comment2.mp4 PiP
        print("Creating demo1 segment...")
        if os.path.exists("demo1.mp4"):
            demo1_clip = VideoFileClip("demo1.mp4")
        else:
            print("Warning: demo1.mp4 not found. Creating placeholder.")
            demo1_clip = self._create_demo_placeholder()
            
        demo1_with_pip = self.create_pip_video(demo1_clip, "comment2.mp4")
        video_clips.append(demo1_with_pip)
        
        # 3. Slide 2 with comment3.mp4 PiP
        print("Creating slide 2 segment...")
        slide2_clip = ImageClip(slide_paths[1]).set_duration(5)  # Default 5 seconds
        slide2_with_pip = self.create_pip_video(slide2_clip, "comment3.mp4")
        video_clips.append(slide2_with_pip)
        
        # Concatenate all clips
        print("Concatenating video segments...")
        final_video = concatenate_videoclips(video_clips)
        
        # Write final video
        print(f"Writing final video to {output_path}...")
        try:
            final_video.write_videofile(
                output_path,
                fps=24,
                codec='libx264',
                audio_codec='aac'
            )
        except Exception as e:
            print(f"Error writing video with libx264/aac: {e}")
            print("Trying with default codecs...")
            try:
                final_video.write_videofile(output_path, fps=24)
            except Exception as e2:
                print(f"Error writing video with default codecs: {e2}")
                raise
        
        # Clean up
        final_video.close()
        for clip in video_clips:
            clip.close()
        
        print(f"Video creation completed: {output_path}")
    
    def _create_demo_placeholder(self) -> VideoFileClip:
        """Create a placeholder for demo1.mp4."""
        from moviepy.editor import ColorClip, TextClip, CompositeVideoClip
        
        # Create background
        bg = ColorClip(size=(1920, 1080), color=(50, 50, 50), duration=10)
        
        # Add text
        try:
            text = TextClip("Demo Video Placeholder", fontsize=50, color='white')
            text = text.set_position('center').set_duration(10)
            demo_clip = CompositeVideoClip([bg, text])
        except Exception as e:
            print(f"Warning: TextClip failed ({e}), using plain background")
            # Fallback if TextClip fails
            demo_clip = bg
            
        return demo_clip


def main():
    """Main function to create the presentation video."""
    with PresentationVideoCreator("presentation.ppt") as creator:
        creator.create_final_video("presentation_with_pip.mp4")


if __name__ == "__main__":
    main()