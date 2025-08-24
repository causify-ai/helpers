#!/usr/bin/env python3

import argparse
import logging
import os
import sys
from pathlib import Path

import helpers.hdbg as hdbg
import helpers.hparser as hparser
import helpers.hsystem as hsystem

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import tempfile
import PyMuPDF as fitz

_LOG = logging.getLogger(__name__)


def extract_slide_images(ppt_path: str, output_dir: str) -> bool:
    """
    Extract embedded images from PowerPoint presentation.
    
    :param ppt_path: path to PowerPoint presentation file
    :param output_dir: directory to save extracted images
    :return: True if extraction successful
    """
    hdbg.dassert(os.path.exists(ppt_path), f"PowerPoint file not found: {ppt_path}")
    _LOG.debug(f"Extracting images from {ppt_path} to {output_dir}")
    presentation = Presentation(ppt_path)
    image_count = 0
    # Process each slide to extract images
    for slide_num, slide in enumerate(presentation.slides, 1):
        for shape_num, shape in enumerate(slide.shapes):
            # Check if shape contains an image
            if hasattr(shape, "image") and shape.image:
                image = shape.image
                # Get image format and data
                image_format = image.ext
                image_data = image.blob
                # Save image to file
                image_filename = f"slide_{slide_num:03d}_image_{shape_num:03d}.{image_format}"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_data)
                image_count += 1
                _LOG.debug(f"Extracted image: {image_filename}")
    _LOG.info(f"Extracted {image_count} images from presentation")
    return True


def extract_slides_as_png(ppt_path: str, output_dir: str) -> bool:
    """
    Extract each slide as a PNG image by converting PPT to PDF first, then PDF to images.
    
    :param ppt_path: path to PowerPoint presentation file
    :param output_dir: directory to save slide images
    :return: True if extraction successful
    """
    hdbg.dassert(os.path.exists(ppt_path), f"PowerPoint file not found: {ppt_path}")
    _LOG.debug(f"Converting slides from {ppt_path} to PNG images")
    # Create temporary file for PDF conversion.
    with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf:
        temp_pdf_path = temp_pdf.name
    # Convert PPT to PDF using LibreOffice.
    cmd = f'libreoffice --headless --convert-to pdf --outdir "{os.path.dirname(temp_pdf_path)}" "{ppt_path}"'
    _LOG.debug(f"Converting PPT to PDF: {cmd}")
    rc, output = hsystem.system_to_string(cmd, abort_on_error=False)
    if rc != 0:
        _LOG.error(f"LibreOffice conversion failed: {output}")
        os.unlink(temp_pdf_path)
        return False
    # Find the actual PDF file created (LibreOffice names it based on input file).
    ppt_basename = os.path.splitext(os.path.basename(ppt_path))[0]
    actual_pdf_path = os.path.join(os.path.dirname(temp_pdf_path), f"{ppt_basename}.pdf")
    if not os.path.exists(actual_pdf_path):
        _LOG.error(f"PDF file not created at expected location: {actual_pdf_path}")
        return False
    # Convert PDF pages to PNG images.
    _LOG.debug(f"Converting PDF to PNG images from {actual_pdf_path}")
    pdf_document = fitz.open(actual_pdf_path)
    slide_count = 0
    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        # Render page as image with high resolution.
        matrix = fitz.Matrix(2.0, 2.0)  # 2x zoom for better quality.
        pix = page.get_pixmap(matrix=matrix)
        # Save as PNG.
        slide_filename = f"slide_{page_num + 1:03d}.png"
        slide_path = os.path.join(output_dir, slide_filename)
        pix.save(slide_path)
        slide_count += 1
        _LOG.debug(f"Saved slide: {slide_filename}")
    pdf_document.close()
    # Clean up temporary files.
    os.unlink(actual_pdf_path)
    if os.path.exists(temp_pdf_path):
        os.unlink(temp_pdf_path)
    _LOG.info(f"Extracted {slide_count} slides as PNG images")
    return True


def extract_notes_text(ppt_path: str, output_dir: str):
    """
    Extract notes from PowerPoint presentation.
    
    :param ppt_path: path to PowerPoint presentation file
    :param output_dir: directory to save extracted notes
    """
    # TODO(ai): Use hdbg.dassert_file_exists()
    hdbg.dassert(os.path.exists(ppt_path), f"PowerPoint file not found: {ppt_path}")
    _LOG.debug(f"Extracting notes from {ppt_path}")
    presentation = Presentation(ppt_path)
    all_notes = []
    # Process each slide to extract notes
    for i, slide in enumerate(presentation.slides, 1):
        slide_notes = []
        # Check if slide has notes and extract them
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            # Extract text from all shapes in notes slide
            for shape in notes_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_notes.append(shape.text.strip())
        # Prepare notes content for this slide
        notes_content = "\n".join(slide_notes) if slide_notes else "No notes"
        # Write individual slide notes file
        slide_notes_path = os.path.join(output_dir, f"slide_{i:03d}_notes.txt")
        with open(slide_notes_path, 'w', encoding='utf-8') as f:
            f.write(f"Slide {i} Notes:\n")
            f.write("=" * 20 + "\n")
            f.write(notes_content)
            f.write("\n")
        all_notes.append(f"Slide {i}:\n{notes_content}\n")
    # Write combined notes file
    all_notes_path = os.path.join(output_dir, "all_notes.txt")
    # TODO(ai): Use hsystem.system() and abort_on_error=True
    with open(all_notes_path, 'w', encoding='utf-8') as f:
        f.write("PowerPoint Presentation Notes\n")
        f.write("=" * 30 + "\n\n")
        f.write("\n".join(all_notes))
    _LOG.info(f"Extracted notes for {len(presentation.slides)} slides")


def extract_text_content(ppt_path: str, output_dir: str):
    """
    Extract all text content from slides.
    
    :param ppt_path: path to PowerPoint presentation file
    :param output_dir: directory to save extracted text content
    """
    hdbg.dassert(os.path.exists(ppt_path), f"PowerPoint file not found: {ppt_path}")
    _LOG.debug(f"Extracting text content from {ppt_path}")
    presentation = Presentation(ppt_path)
    all_text = []
    # Process each slide to extract text content
    for i, slide in enumerate(presentation.slides, 1):
        slide_text = []
        # Extract text from all shapes in slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            # Handle table shapes specifically
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            slide_text.append(cell.text.strip())
        # Prepare text content for this slide
        text_content = "\n".join(slide_text) if slide_text else "No text content"
        # Write individual slide text file
        slide_text_path = os.path.join(output_dir, f"slide_{i:03d}_text.txt")
        with open(slide_text_path, 'w', encoding='utf-8') as f:
            f.write(f"Slide {i} Text Content:\n")
            f.write("=" * 25 + "\n")
            f.write(text_content)
            f.write("\n")
        all_text.append(f"Slide {i}:\n{text_content}\n")
    # Write combined text content file
    all_text_path = os.path.join(output_dir, "all_text_content.txt")
    with open(all_text_path, 'w', encoding='utf-8') as f:
        f.write("PowerPoint Presentation Text Content\n")
        f.write("=" * 35 + "\n\n")
        f.write("\n".join(all_text))
    _LOG.info(f"Extracted text content for {len(presentation.slides)} slides")


def _check_platform() -> None:
    """
    Check that the script is running on Linux or MacOS.
    
    :raises: SystemExit if running on unsupported platform
    """
    platform = sys.platform.lower()
    if not (platform.startswith('linux') or platform.startswith('darwin')):
        _LOG.error(f"Unsupported platform: {platform}. This script only works on Linux and MacOS.")
        sys.exit(1)
    _LOG.debug(f"Running on supported platform: {platform}")


def main():
    """
    Main function to extract slides, notes, and text content from PowerPoint presentations.
    """
    # Check platform compatibility
    _check_platform()
    # TODO(ai): Move this in a _parse() function like in the script_template.py
    # Parse command line arguments
    parser = argparse.ArgumentParser(
        description="Extract images, text content and notes from PowerPoint presentations"
    )
    hparser.add_verbosity_arg(parser)
    parser.add_argument("ppt_file", help="Path to PowerPoint presentation file")
    parser.add_argument(
        "-o", "--output-dir", 
        help="Output directory (default: presentation_name_extracted)"
    )
    parser.add_argument(
        "--extract-images", 
        action="store_true",
        help="Extract embedded images from slides"
    )
    parser.add_argument(
        "--extract-slides", 
        action="store_true",
        help="Extract each slide as a PNG image"
    )
    #
    args = parser.parse_args()
    hdbg.init_logger(verbosity=args.log_level, use_exec_path=True)
    # Validate input file
    ppt_path = args.ppt_file
    hdbg.dassert(os.path.exists(ppt_path), f"PowerPoint file not found: {ppt_path}")
    # Determine output directory
    if args.output_dir:
        output_dir = args.output_dir
    else:
        ppt_name = Path(ppt_path).stem
        output_dir = f"{ppt_name}_extracted"
    # TODO(ai): Use hio.create_dir()
    os.makedirs(output_dir, exist_ok=True)
    #
    _LOG.info(f"Extracting presentation: {ppt_path}")
    _LOG.info(f"Output directory: {output_dir}")
    # Extract embedded images if requested
    if args.extract_images:
        _LOG.info("Extracting embedded images...")
        success = extract_slide_images(ppt_path, output_dir)
        if success:
            _LOG.info("Successfully extracted embedded images")
        else:
            _LOG.warning("No images found or extraction failed")
    # Extract notes from presentation
    _LOG.info("Extracting notes...")
    extract_notes_text(ppt_path, output_dir)
    # Extract text content from slides
    _LOG.info("Extracting text content...")
    extract_text_content(ppt_path, output_dir)
    _LOG.info(f"Extraction complete! Check the '{output_dir}' directory.")


if __name__ == "__main__":
    main()